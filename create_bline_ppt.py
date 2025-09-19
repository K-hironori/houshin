from pptx import Presentation
from pptx.util import Pt

prs = Presentation()

TITLE_FONT = "Noto Sans JP"
BODY_FONT = "Noto Sans JP"

def set_title_shape(shape, text):
    shape.text = text
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = TITLE_FONT
            run.font.size = Pt(44)

def set_body_text(shape, lines):
    text_frame = shape.text_frame
    text_frame.clear()
    for idx, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        paragraph.text = line
        paragraph.level = 0
        for run in paragraph.runs:
            run.font.name = BODY_FONT
            run.font.size = Pt(24)

# Title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
set_title_shape(slide.shapes.title, "Bライン溶接工程 改善提案")
subtitle = slide.placeholders[1]
subtitle.text = "特型工数入力の標準化と省力化"
for paragraph in subtitle.text_frame.paragraphs:
    for run in paragraph.runs:
        run.font.name = BODY_FONT
        run.font.size = Pt(28)

# Overview
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_title_shape(slide.shapes.title, "概要")
body = slide.shapes.placeholders[1]
set_body_text(body, [
    "2018年稼働開始以降、暫定運用のまま特型工数を管理",
    "2024年1月時点で未設定品番が2,600品番",
    "目標：標準化と省力化により未設定ゼロを実現"
])

# KPI
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_title_shape(slide.shapes.title, "主要KPI")
body = slide.shapes.placeholders[1]
set_body_text(body, [
    "工数未設定特型品番：2,600品番",
    "削減見込み工数：1.0時間/日",
    "創出時間は現場改善・育成へ再投資"
])

# Problem 1
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_title_shape(slide.shapes.title, "課題① 工数算出の属人化")
body = slide.shapes.placeholders[1]
set_body_text(body, [
    "図面を毎回目視し、個人の経験と勘に依存",
    "判断基準が曖昧で工数にバラつきが発生",
    "ノウハウが個人に滞留し継承が困難",
    "消費工数：約0.5時間/日"
])

# Problem 2
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_title_shape(slide.shapes.title, "課題② 実績入力の非効率")
body = slide.shapes.placeholders[1]
set_body_text(body, [
    "指示票を参照しながらシステムへ手入力",
    "入力ミス・記録漏れなどヒューマンエラーリスク",
    "データ信頼性が低く分析が困難",
    "消費工数：約0.5時間/日"
])

# Solution direction
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_title_shape(slide.shapes.title, "解決の方向性")
body = slide.shapes.placeholders[1]
set_body_text(body, [
    "属人化から仕組み化へ転換",
    "マスター登録と自動転記で入力作業を削減",
    "誰が担当しても同じ工数で運用できる体制を構築"
])

# Steps
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_title_shape(slide.shapes.title, "改善ステップ")
body = slide.shapes.placeholders[1]
set_body_text(body, [
    "STEP1 ルール化：標準工数と判断基準を定義",
    "STEP2 マスター登録：品番ごとに標準工数を登録",
    "STEP3 自動化：実績に応じてマスター工数を自動転記"
])

# Plan
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_title_shape(slide.shapes.title, "削減計画")
body = slide.shapes.placeholders[1]
set_body_text(body, [
    "開始：2024年11月",
    "ペース：月176品番削減",
    "期間：24か月で未設定ゼロ（2026年10月目標）",
    "ポイント：新規発生分を処理しつつ未処理分を段階的に解消"
])

# Quantitative effects
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_title_shape(slide.shapes.title, "定量効果")
body = slide.shapes.placeholders[1]
set_body_text(body, [
    "組長工数を1.0時間/日削減",
    "月20日稼働×12か月＝年間240時間創出",
    "創出時間を品質改善・育成などコア業務へ配分"
])

# Qualitative effects
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_title_shape(slide.shapes.title, "定性効果")
body = slide.shapes.placeholders[1]
set_body_text(body, [
    "業務の標準化：公平で透明性の高い工数管理",
    "データ信頼性の向上：ヒューマンエラーを排除",
    "将来の生産計画・原価計算の精度向上に貢献"
])

# Closing
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_title_shape(slide.shapes.title, "期待される成果")
body = slide.shapes.placeholders[1]
set_body_text(body, [
    "標準化と省力化による持続可能な生産体制",
    "工数データの正確性向上で戦略的な改善が可能",
    "創出時間をコア業務へ再投資し競争力を強化"
])

prs.save("Bライン溶接工程_改善提案.pptx")
